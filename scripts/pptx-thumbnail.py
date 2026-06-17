"""
Gera thumbnail de cada slide do PPTX para QA visual.
Usa python-pptx para extrair informações e PIL para montar imagem de preview simplificada.
"""
import sys, os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import base64, io

pptx_path = sys.argv[1] if len(sys.argv) > 1 else "treinamento-app-brigadas-2026.pptx"
out_dir = sys.argv[2] if len(sys.argv) > 2 else "tmp/slide-thumbs"
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(pptx_path)
slide_w = prs.slide_width
slide_h = prs.slide_height

THUMB_W = 800
THUMB_H = int(THUMB_W * slide_h / slide_w)

print(f"Gerando {len(prs.slides)} thumbnails ({THUMB_W}x{THUMB_H}px)...")

for i, slide in enumerate(prs.slides):
    img = Image.new('RGB', (THUMB_W, THUMB_H), '#F0F0F0')
    draw = ImageDraw.Draw(img)

    # Render background
    bg = slide.background
    fill = bg.fill
    if fill.type is not None:
        try:
            color = fill.fore_color.rgb
            img = Image.new('RGB', (THUMB_W, THUMB_H), f'#{color}')
            draw = ImageDraw.Draw(img)
        except Exception:
            pass

    # Render shapes (simplified)
    for shape in slide.shapes:
        try:
            # Get position/size in pixels
            l = int(shape.left / slide_w * THUMB_W)
            t = int(shape.top / slide_h * THUMB_H)
            w = int(shape.width / slide_w * THUMB_W)
            h = int(shape.height / slide_h * THUMB_H)

            # Draw shape fill if available
            if hasattr(shape, 'fill'):
                try:
                    fill = shape.fill
                    if fill.type is not None and str(fill.type) == 'SOLID (1)':
                        color = f'#{fill.fore_color.rgb}'
                        draw.rectangle([l, t, l+w, t+h], fill=color)
                except Exception:
                    draw.rectangle([l, t, l+w, t+h], fill='#E0E0E0')

            # Draw embedded images
            if shape.shape_type == 13:  # PICTURE
                try:
                    img_bytes = shape.image.blob
                    sub_img = Image.open(io.BytesIO(img_bytes)).convert('RGB')
                    sub_img = sub_img.resize((max(1,w), max(1,h)), Image.LANCZOS)
                    img.paste(sub_img, (l, t))
                except Exception:
                    draw.rectangle([l, t, l+w, t+h], fill='#CCCCCC')
                    draw.text((l+4, t+4), '[img]', fill='#666666')

            # Draw text
            if shape.has_text_frame:
                text = shape.text_frame.text[:80]
                if text.strip():
                    try:
                        font = ImageFont.truetype("arial.ttf", max(8, min(18, int(h * 0.4))))
                    except Exception:
                        font = ImageFont.load_default()
                    try:
                        txt_color = f'#{shape.text_frame.paragraphs[0].runs[0].font.color.rgb}'
                    except Exception:
                        txt_color = '#333333'
                    draw.text((l+2, t+2), text, fill=txt_color, font=font)

        except Exception:
            pass

    # Save with slide number label
    draw.rectangle([0, THUMB_H-24, THUMB_W, THUMB_H], fill='#000000AA')
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except Exception:
        font = ImageFont.load_default()
    draw.text((8, THUMB_H-20), f'Slide {i+1}', fill='#FFFFFF', font=font)

    out_path = os.path.join(out_dir, f'slide-{i+1:02d}.jpg')
    img.save(out_path, 'JPEG', quality=85)
    print(f'  ✓ slide-{i+1:02d}.jpg')

print(f'\nThumbnails em: {out_dir}')
